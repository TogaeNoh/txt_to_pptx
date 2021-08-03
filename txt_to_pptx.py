# import required module
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import ImageFont
import datetime


def main():
    text_filename = r'C:\Users\muric\Desktop\Togae\Writing\210629첫점심.txt'
    para_list = file_to_paragraph(text_filename)
    words = para_to_words(para_list)

    font_size = 18
    font_name = '경기천년바탕OTF_Regular.otf'
    widths_w = width_words(words, font_size, font_name)

    text_width_px = 360
    line_num_per_slide = 13

    line_list = words_to_line(words, widths_w, text_width_px)
    slides_list = line_to_slide(line_list, line_num_per_slide)

    ppt = create_slides(slides_list)

    # save file
    filefolder = r'C:\Users\muric\Desktop\Togae\writingPPT'
    today = datetime.datetime.today()
    filename = os.path.join(filefolder, f'{today:%Y%m%d}.pptx')

    ppt.save(filename)


def width_words(words, font_size, font_name):
    num_para = len(words)

    width_words_lst = []
    for i in range(num_para):
        width_words_lst.append([])

    for i in range(num_para):
        for j in range(len(words[i])):
            width_words_lst[i].append(get_pil_text_size(words[i][j], font_size, font_name)[0])

    return width_words_lst


def create_word_list(paragraph_list):
    word_list = []
    for i in range(len(paragraph_list)):
        for word in paragraph_list[i].split():
            word_list[i].append(word)

    return word_list


def get_pil_text_size(text, font_size, font_name):
    font = ImageFont.truetype(font_name, font_size)
    size = font.getsize(text)
    return size


def file_to_paragraph(filename):
    # Extract words from text file
    paragraph = []

    # create a paragraph
    with open(filename, encoding='utf-8')  as f:
        for line in f:
            paragraph.append('\t' + line[:len(line)-1]) # remove '\n'

    return paragraph


def para_to_words(paragraph):
    words = []
    for i in range(len(paragraph)):
        words.append([])

    for i in range(len(paragraph)):
        for word in paragraph[i].split():
            #print(word)
            words[i].append(word + ' ')

    for i in range(len(words)):
        words[i][0] = '\t'+ words[i][0]

    return words


def words_to_line(words, widths_w, text_width_px):
    num_para = len(words)
    line_list = []

    for i in range(num_para):
        line_list.append([])

    for i in range(num_para):
        while words[i]:
            for j in range(len(words[i])):
                if sum(widths_w[i][:j]) >= text_width_px :
                    line_str = "".join(words[i][0:j-1])
                    line_str += '\n'
                    line_list[i].append(line_str)
                    words[i] = words[i][j-1:]
                    widths_w[i]= widths_w[i][j-1:]
                    break
                elif j == (len(words[i])-1): #last line of each paragraph
                    if sum(widths_w[i]) > text_width_px:
                        line_str = "".join(words[i][:-1])
                        line_str += '\n'
                        line_list[i].append(line_str)
                        line_list[i].append( words[i][-1:][0] + '\n' )
                        words[i] = ''
                    else:
                        line_str = "".join(words[i][:])
                        line_list[i].append(line_str)
                        line_str += '\n'
                        words[i] = ''
                    break

    flat_list = [item for sublist in line_list for item in sublist]

    return flat_list


def line_to_slide(line_list, line_num_per_slide):
    slide_num = len(line_list)//line_num_per_slide + 1
    slide_list=[]
    for i in range(slide_num):
        if i == (slide_num-1): #last slide
            slide_str= "".join(line_list[ i*line_num_per_slide : ])
            slide_list.append(slide_str)
        else:
            slide_str = "".join(line_list[ i*line_num_per_slide : ((i+1)*line_num_per_slide)])
            slide_list.append(slide_str)

    return slide_list


def create_slides(slide_list):
    # Creating Object
    ppt = Presentation()
    ppt.slide_width = Cm(15)
    ppt.slide_height = Cm(15)

    slides_number = len(slide_list)

    # To create blank slide layout
    blank_slide_layout = {}
    slide = {}
    txBox = {}
    tf = {}
    p = {}
    run = {}
    font = {}

    # For adjusting the  Margins in inches
    left = Cm(.5)
    top = Cm(.01)
    width = Cm(14.5)
    height = Cm(15)

    for x in range(slides_number):
        blank_slide_layout[x] = ppt.slide_layouts[6]
        slide[x] = ppt.slides.add_slide(blank_slide_layout[x])
        txBox[x] = slide[x].shapes.add_textbox(left, top, width, height)

        tf[x] = txBox[x].text_frame
        tf[x].margin_left = Cm(0.72)
        tf[x].margin_right = Cm(0.000001)

        tf[x].margin_top = Cm(1)
        tf[x].margin_bottom = Cm(1)
        tf[x].word_wrap = False
        tf[x].auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        p[x] = tf[x].paragraphs[0]
        p[x].line_spacing = 1.3
        run[x] = p[x].add_run()

        run[x].text = slide_list[x]
        font[x] = run[x].font
        font[x].name = '경기천년바탕OTF Regular'
        font[x].size = Pt(18)

    return ppt

if __name__ == '__main__':
    main()